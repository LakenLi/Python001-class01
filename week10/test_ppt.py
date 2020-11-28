# encoding: utf-8
'''
@author: horan
@license: (C) Copyright 2020-
@contact: laken_phil@163.com
@file: test_ppt.py
@time: 2020/11/22 上午10:48
@desc:
'''

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.util import Inches,Pt,Cm
from pptx.chart.data import ChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK, XL_DATA_LABEL_POSITION


# 使用自定义模版
prs = Presentation('/Users/horan/horan_work_for_git/Python001-class01/week10/caicai_demo.pptx')

# 处理第一张
slide_01 = prs.slides[0]
# 替换第一张标题
slide_01.shapes[0].text_frame.text = '12·29 Shopee专场·XX站'
slide_01.shapes[3].text_frame.text = '品牌中心 | 2020.12.29'

# 第二、三、四张幻灯片不处理
# 开始处理第五张
slide_05 = prs.slides[4]
slide_05.shapes[23].text_frame.text = '8天'                                                  # 替换推广时长
slide_05.shapes[27].text_frame.text = '11.22-11.29'                                           # 替换推广起止日期
slide_05.shapes[24].text_frame.text = '2344UV'                                               # 替换专题流量
slide_05.shapes[28].text_frame.text = '浏览次数3333'                                          # 替换浏览次数
slide_05.shapes[25].text_frame.text = '300人'                                                # 替换报名人数
slide_05.shapes[29].text_frame.text = '转化率 ' + str(round(300 / 2344, 4) * 100) + '%'       # 替换转化率
slide_05.shapes[26].text_frame.text = '211人'                                                # 替换到场人数
slide_05.shapes[30].text_frame.text = '到场率' + str(round(211 / 300, 4) * 100) + '%'        # 替换到场率
slide_05.shapes[1].text_frame.text = '注：活动页面从11.22开始上线加入“春雨计划”系列活动推广。'     # 替换注释里的日期


# 开始处理第六张
slide_06 = prs.slides[5]
# 替换内容
slide_06.shapes[1].text_frame.text = '11.22-11.25'     # 第一波替换
slide_06.shapes[4].text_frame.text = '11.26-11.28'      # 第二波替换
slide_06.shapes[6].text_frame.text = '"此次活动推广时间为?天，整体推广内容上分为2波（主题议程→嘉宾（招商经理）），整体推广节奏及报名效果达到预期。"'      # 推广活动描述替换
slide_06.shapes[10].text_frame.text = '报名：70人'      # 嘉宾内容替换
slide_06.shapes[13].text_frame.text = '报名：88人'     # 主题议程替换
slide_06.shapes[14].text_frame.text = '11.28'     # 活动期替换
slide_06.shapes[15].text_frame.text = '到场：77人'      # 签到+转化 替换


# 开始处理第7张
slide_07 = prs.slides[6]
# 替换内容
slide_07.shapes[12].text_frame.text = '注：未知来源有？人。'
slide_07.shapes[20].text_frame.text = '33人'               # 弹窗箭头指向人数替换
slide_07.shapes[22].text_frame.text = '报名：11人'          # 会员广告位+开店banner 报名人数替换
slide_07.shapes[23].text_frame.text = '报名：3人'          # 卖家活动页面 报名人数替换
slide_07.shapes[24].text_frame.text = '报名：750人'          # 硬广位+果园+弹窗 报名人数替换
slide_07.shapes[25].text_frame.text = '报名：11人'          # 海报报名人数替换
slide_07.shapes[26].text_frame.text = '报名：110人'          # 自然流量 报名人数替换


# 忽略第八、九张

# 开始处理第十张
slide_10 = prs.slides[9]
shapes_10 = slide_10.shapes

# 定义图表数据-------------------------------------------------
x = ['11月22日', '11月23日', '11月24日', '11月25日', '11月26日', '11月27日']
y = [88, 89, 77, 998, 34, 34]
y2 = [23, 145, 356, 245, 45, 555]
y3 = [234, 233, 123, 44, 55, 66]

chart_data = ChartData()
chart_data.categories = x
chart_data.add_series(name='浏览次数', values=y)
chart_data.add_series(name='浏览人数', values=y2)
chart_data.add_series(name='报名人数', values=y3)

# 添加图表-------------------------------------------------
left, top, width, height = Inches(1), Inches(1.5), Inches(12), Inches(6)
graphic_frame = shapes_10.add_chart(chart_type=XL_CHART_TYPE.LINE_MARKERS,  # 图表类型
                                 x=left, y=top,    # 图表区的位置
                                 cx=width, cy=height,    # 图表的宽和高
                                 chart_data=chart_data)

slide_10.shapes[0].text_frame.text = '流量及报名趋势'
slide_10.shapes[2].text_frame.text = '注：数据统计范围为11.22-11.27'

# 开始处理第11张
slide_11 = prs.slides[10]
slide_11.shapes[27].text_frame.text = '2333'                                   # 替换活动页数据
slide_11.shapes[31].text_frame.text = '455'                                    # 替换报名页数据
slide_11.shapes[32].text_frame.text = '321'                                    # 替换报名成功页数据
slide_11.shapes[28].text_frame.text = str(round(455 / 2333, 4) * 100) + '%'    # 报名页/活动页
slide_11.shapes[29].text_frame.text = str(round(321 / 455, 4) * 100) + '%'     # 报名成功页/报名页
slide_11.shapes[30].text_frame.text = str(round(321 / 2333, 4) * 100) + '%'    # 报名成功页/活动页
slide_11.shapes[21].text_frame.text = '注：数据统计范围为11.22-11.28。'           # 替换注释日期

# 开始处理第12张
slide_12 = prs.slides[11]

shapes_12 = slide_12.shapes
# 定义图表数据-------------------------------------------------
x = ['弹窗', '首页banner', '信息流B01', '海报1', '首页轮播', 'APP轮播', '果园banner', '海报2', '开店banner', '短信']
y = [729, 201, 162, 114, 96, 86, 75, 54, 53, 30]

chart_data = ChartData()
chart_data.categories = x
chart_data.add_series(name='活动页流量来源', values=y)

# 添加图表-------------------------------------------------
left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(6)
graphic_frame_12 = shapes_12.add_chart(chart_type=XL_CHART_TYPE.PIE,  # 图表类型
                                 x=left, y=top,    # 图表区的位置
                                 cx=width, cy=height,    # 图表的宽和高
                                 chart_data=chart_data)

chart = graphic_frame_12.chart
chart.chart_style = 26
plot = chart.plots[0]
# 设置数据标签
plot.has_data_labels = True    # 显示数据标签
data_labels = plot.data_labels    # 获取数据标签控制类
data_labels.show_category_name = True    # 是否显示类别名称
data_labels.show_value = False    # 是否显示值
data_labels.show_percentage = True    # 是否显示百分比
data_labels.number_format = '0.0%'    # 标签的数字格式
data_labels.position = XL_DATA_LABEL_POSITION.INSIDE_END    # 标签位置
chart.font.name = '微软雅黑'
chart.font.size = Pt(9)

# 设置图表标题
chart.has_title = True    # 显示标题
para = chart.chart_title.text_frame.add_paragraph()
para.text = '活动页流量来源分布图'    # 标题内容
para.font.size = Pt(16)    # 字体大小

# 添加表格
rows, cols = 11, 2  # 设定11行 2列
left = Cm(22)
top = Cm(5)
width = Cm(18)
height = Cm(12)
table_12 = shapes_12.add_table(rows, cols, left, top, width, height).table

table_12.columns[0].width=Cm(5) # 列宽
table_12.columns[1].width=Cm(5)
table_12.rows[0].heigth=Cm(6)# 行宽
data=[
    ['渠道', '浏览用户'],
    ['弹窗', 729],
    ['首页banner', 201],
    ['信息流B01', 162],
    ['海报1', 114],
    ['首页轮播', 96],
    ['APP轮播', 86],
    ['果园banner', 75],
    ['海报2', 54],
    ['开店banner', 53],
    ['短信', 30]
]

for row in range(rows): # 循环行
    for col in range(cols): # 循环列
        table_12.cell(row, col).text_frame.clear()
        new_cell = table_12.cell(row, col).text_frame.paragraphs[0]
        new_cell.text = str(data[row][col])
        new_cell.font.name = '微软雅黑'
        new_cell.font.size = Pt(12)
        new_cell.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER  # 居中


## ==================================================================================================================
# 开始处理第13张
slide_13 = prs.slides[12]

shapes_13 = slide_13.shapes
# 定义图表数据-------------------------------------------------
x = ['弹窗', '海报1', '首页banner', '海报2', '开店banner', '果园banner', 'APP轮播', '信息流B01', '其他活动', '短信', '首页轮播', '直播banner', '亚马逊B04', '未知']
y = [35, 16, 14, 11, 8, 8, 8, 7, 3, 3, 2, 1, 17]

chart_data = ChartData()
chart_data.categories = x
chart_data.add_series(name='活动页流量来源', values=y)

# 添加图表-------------------------------------------------
left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(6)
graphic_frame_13 = shapes_13.add_chart(chart_type=XL_CHART_TYPE.PIE,  # 图表类型
                                 x=left, y=top,    # 图表区的位置
                                 cx=width, cy=height,    # 图表的宽和高
                                 chart_data=chart_data)

chart = graphic_frame_13.chart
chart.chart_style = 26
plot = chart.plots[0]
# 设置数据标签
plot.has_data_labels = True    # 显示数据标签
data_labels = plot.data_labels    # 获取数据标签控制类
data_labels.show_category_name = True    # 是否显示类别名称
data_labels.show_value = False    # 是否显示值
data_labels.show_percentage = True    # 是否显示百分比
data_labels.number_format = '0.0%'    # 标签的数字格式
data_labels.position = XL_DATA_LABEL_POSITION.INSIDE_END    # 标签位置
chart.font.name = '微软雅黑'
chart.font.size = Pt(8)

# 设置图表标题
chart.has_title = True    # 显示标题
para = chart.chart_title.text_frame.add_paragraph()
para.text = '观众报名来源分布图'    # 标题内容
para.font.size = Pt(16)    # 字体大小

# 添加表格
rows, cols = 15, 2  # 设定15行 2列
left = Cm(22)
top = Cm(5)
width = Cm(18)
height = Cm(12)
table_13 = shapes_13.add_table(rows, cols, left, top, width, height).table

table_13.columns[0].width=Cm(5) # 列宽
table_13.columns[1].width=Cm(5)
table_13.rows[0].heigth=Cm(6)# 行宽
data=[
    ['渠道', '流量'],
    ['弹窗', 35],
    ['海报1', 16],
    ['首页banner', 14],
    ['海报2', 11],
    ['开店banner', 8],
    ['果园banner', 8],
    ['APP轮播', 8],
    ['信息流B01', 7],
    ['其他活动', 3],
    ['短信', 30],
    ['首页轮播', 96],
    ['直播banner', 1],
    ['亚马逊B04', 1],
    ['未知', 17]
]

for row in range(rows): # 循环行
    for col in range(cols): # 循环列
        table_13.cell(row, col).text_frame.clear()
        new_cell = table_13.cell(row, col).text_frame.paragraphs[0]
        new_cell.text = str(data[row][col])
        new_cell.font.name = '微软雅黑'
        new_cell.font.size = Pt(12)
        new_cell.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER  # 居中


# ==============================================================================================
# 开始处理第14张
slide_14 = prs.slides[13]

shapes_14 = slide_14.shapes
# 添加表格1流量
rows, cols = 11, 2  # 设定11行 2列
left = Cm(3)
top = Cm(5)
width = Cm(18)
height = Cm(11)
table_14_01 = shapes_14.add_table(rows, cols, left, top, width, height).table

table_14_01.columns[0].width=Cm(4) # 列宽
table_14_01.columns[1].width=Cm(4)
table_14_01.rows[0].heigth=Cm(5)# 行宽
# data_01 为第一个流量表格，展示渠道与流量
data_01=[
    ['渠道', '流量'],
    ['弹窗', 729],
    ['首页banner', 201],
    ['信息流B01', 162],
    ['海报1', 114],
    ['首页轮播', 96],
    ['APP轮播', 86],
    ['果园banner', 75],
    ['海报2', 54],
    ['开店banner', 53],
    ['短信', 30]
]

for row in range(rows): # 循环行
    for col in range(cols): # 循环列
        table_14_01.cell(row, col).text_frame.clear()
        new_cell = table_14_01.cell(row, col).text_frame.paragraphs[0]
        new_cell.text = str(data_01[row][col])
        new_cell.font.name = '微软雅黑'
        new_cell.font.size = Pt(12)
        new_cell.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER  # 居中


# 添加表格2报名
rows, cols = 11, 2  # 设定11行 2列
left = Cm(13)
top = Cm(5)
width = Cm(18)
height = Cm(11)
table_14_02 = shapes_14.add_table(rows, cols, left, top, width, height).table

table_14_02.columns[0].width=Cm(4) # 列宽
table_14_02.columns[1].width=Cm(4)
table_14_02.rows[0].heigth=Cm(5)# 行宽
# data_02 为第一个流量表格，展示渠道与报名
data_02=[
    ['渠道', '报名'],
    ['弹窗', 35],
    ['首页banner', 14],
    ['信息流B01', 7],
    ['海报1', 16],
    ['首页轮播', 2],
    ['APP轮播', 8],
    ['果园banner', 8],
    ['海报2', 11],
    ['开店banner', 8],
    ['短信', 3]
]

for row in range(rows): # 循环行
    for col in range(cols): # 循环列
        table_14_02.cell(row, col).text_frame.clear()
        new_cell = table_14_02.cell(row, col).text_frame.paragraphs[0]
        new_cell.text = str(data_02[row][col])
        new_cell.font.name = '微软雅黑'
        new_cell.font.size = Pt(12)
        new_cell.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER  # 居中


# 添加表格3转化率
rows, cols = 11, 2  # 设定11行 2列
left = Cm(23)
top = Cm(5)
width = Cm(18)
height = Cm(11)
table_14_03 = shapes_14.add_table(rows, cols, left, top, width, height).table

table_14_03.columns[0].width=Cm(4) # 列宽
table_14_03.columns[1].width=Cm(4)
table_14_03.rows[0].heigth=Cm(5)# 行宽

for row in range(rows): # 循环行
    for col in range(cols): # 循环列
        table_14_03.cell(row, col).text_frame.clear()
        new_cell = table_14_03.cell(row, col).text_frame.paragraphs[0]
        if row == 0 & col == 0 :
            new_cell.text = '渠道'
        elif row == 0 & col == 1:
            new_cell.text = '转化率'
        else:
            if col == 0:
                new_cell.text = data_01[row][col]
            else:
                new_cell.text = str(round(data_02[row][col] / data_01[row][col], 4) * 100) + '%'
        new_cell.font.name = '微软雅黑'
        new_cell.font.size = Pt(12)
        new_cell.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER  # 居中


i=0

# print(slide_10.shapes.title.text)

# for slide in slide_11.shapes:
#
#     if slide.has_text_frame:
#         # print('第: {0} 个，文本框内容: {1}'.format(i, slide.text_frame.text))
#
#     i += 1
#
# print(i)
prs.save('test_00411.pptx')