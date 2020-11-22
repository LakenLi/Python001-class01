# encoding: utf-8
'''
@author: horan
@license: (C) Copyright 2020-
@contact: laken_phil@163.com
@file: test_2.py
@time: 2020/11/22 下午7:57
@desc:
'''
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK, XL_DATA_LABEL_POSITION

# 创建幻灯片-------------------------------------------------
prs = Presentation()    # 初始化 ppt 文档
title_only_slide_layout = prs.slide_layouts[5]    # 选择空白幻灯片版式
slide = prs.slides.add_slide(title_only_slide_layout)    # 添加一页空白幻灯片
shapes = slide.shapes

shapes.title.text = '折线图'

# 定义图表数据-------------------------------------------------
x = ['11月22日', '11月23日', '11月24日', '11月25日', '11月26日', '11月27日']
y = [8888, 8899, 7788, 9988, 34, 34]
y2 = [23, 4545, 6567, 245, 45, 555]

chart_data = ChartData()
chart_data.categories = x
chart_data.add_series(name='浏览次数', values=y)
chart_data.add_series(name='浏览人数', values=y2)

# 添加图表-------------------------------------------------
left, top, width, height = Inches(1), Inches(1.5), Inches(12), Inches(6)
graphic_frame = shapes.add_chart(chart_type=XL_CHART_TYPE.LINE_MARKERS,  # 图表类型
                                 x=left, y=top,    # 图表区的位置
                                 cx=width, cy=height,    # 图表的宽和高
                                 chart_data=chart_data)
# 保存 ppt 文档
prs.save('test12211.pptx')