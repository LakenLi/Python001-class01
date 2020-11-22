# encoding: utf-8
'''
@author: horan
@license: (C) Copyright 2020-
@contact: laken_phil@163.com
@file: test4.py
@time: 2020/11/22 下午10:07
@desc:
'''
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK, XL_DATA_LABEL_POSITION

# 创建幻灯片-------------------------------------------------
prs = Presentation()    # 初始化 ppt 文档
title_only_slide_layout = prs.slide_layouts[5]    # 选择空白幻灯片版式
slide = prs.slides.add_slide(title_only_slide_layout)    # 添加一页空白幻灯片
shapes = slide.shapes

shapes.title.text = '饼图'

# 定义图表数据-------------------------------------------------
x = ['Q1', 'Q2', 'Q3', 'Q4']
y = [8888, 8899, 7788, 9988]

chart_data = ChartData()
chart_data.categories = x
chart_data.add_series(name='销量', values=y)

# 添加图表-------------------------------------------------
left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(6)
graphic_frame = shapes.add_chart(chart_type=XL_CHART_TYPE.PIE,  # 图表类型
                                 x=left, y=top,    # 图表区的位置
                                 cx=width, cy=height,    # 图表的宽和高
                                 chart_data=chart_data)

chart = graphic_frame.chart
plot = chart.plots[0]
# 设置数据标签
plot.has_data_labels = True    # 显示数据标签
data_labels = plot.data_labels    # 获取数据标签控制类
data_labels.show_category_name = True    # 是否显示类别名称
data_labels.show_value = False    # 是否显示值
data_labels.show_percentage = True    # 是否显示百分比
data_labels.number_format = '0.0%'    # 标签的数字格式
data_labels.position = XL_DATA_LABEL_POSITION.INSIDE_END    # 标签位置
data_labels.font.name = 'Arial'
data_labels.font.size = Pt(14)

# 设置图表标题
chart.has_title = True    # 显示标题
para = chart.chart_title.text_frame.add_paragraph()
para.text = '销量占比'    # 标题内容
para.font.size = Pt(16)    # 字体大小

# 保存 ppt 文档
prs.save('test33.pptx')